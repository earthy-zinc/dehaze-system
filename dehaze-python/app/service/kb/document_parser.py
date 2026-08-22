"""AI 知识库文档解析器：将各类文件解析为纯文本 + 结构化 blocks。

纯算法组件，不依赖数据库。解析策略对齐《后端实现-文档管理.md》§3。
外部解析库（pymupdf/python-docx/python-pptx/openpyxl/rapidocr）全部延迟导入，
未安装时抛出明确业务异常。
"""

import io
from dataclasses import dataclass, field

from app.core.code import ResultCode
from app.core.exceptions import BusinessException

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".md", ".txt", ".jpg", ".png"}
_IMAGE_EXTENSIONS = {".jpg", ".png"}

# 网页正文抽取时去除的噪声标签
_NOISE_TAGS = ["script", "style", "nav", "header", "footer", "aside", "noscript"]


@dataclass
class ParsedDocument:
    content: str
    raw_content: str
    blocks: list[dict] = field(default_factory=list)


def parse_document(file_bytes: bytes, filename: str, strategy: str = "auto") -> ParsedDocument:
    """统一解析入口：按文件格式与解析策略将字节流解析为文本。

    参数：
        file_bytes: 文件二进制内容
        filename: 文件名（用于识别扩展名）
        strategy: auto / fast / hi_res / ocr_only
    返回：
        ParsedDocument（content 纯文本、raw_content 原始结构化文本、blocks 分块列表）
    """
    ext = _extension(filename)
    if ext not in SUPPORTED_EXTENSIONS:
        raise BusinessException(ResultCode.USER_UPLOAD_FILE_TYPE_NOT_MATCH, "不支持的文件格式")
    strategy = (strategy or "auto").lower()

    # 空文件：返回空 ParsedDocument（completed 而非 failed，分块数 0）
    if not file_bytes:
        return ParsedDocument(content="", raw_content="", blocks=[])
    try:
        if ext in _IMAGE_EXTENSIONS:
            blocks = _parse_image(file_bytes)
        elif ext == ".pdf":
            blocks = _parse_pdf(file_bytes, strategy)
        else:
            blocks = _parse_office(file_bytes, ext)
    except BusinessException:
        raise
    except Exception as exc:
        # 底层解析库异常（坏文件/伪装文件）统一转为业务异常，避免 500
        raise BusinessException(ResultCode.BUSINESS_ERROR, f"文档解析失败: {exc}") from exc
    return _build_parsed(blocks)


def parse_html(html: str) -> str:
    """网页正文抽取：去除 script/style/nav 等噪声标签后输出纯文本。"""
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(_NOISE_TAGS):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _extension(filename: str) -> str:
    base = filename or ""
    if "." not in base:
        return ""
    return "." + base.rsplit(".", 1)[-1].lower()


def _build_parsed(blocks: list[dict]) -> ParsedDocument:
    content = "\n\n".join(block["text"] for block in blocks if block.get("text"))
    return ParsedDocument(content=content, raw_content=content, blocks=blocks)


def _blocks_has_text(blocks: list[dict]) -> bool:
    return any(block.get("text", "").strip() for block in blocks)


# ---------------------------------------------------------------------------
# 图片解析（仅 OCR）
# ---------------------------------------------------------------------------
def _parse_image(file_bytes: bytes) -> list[dict]:
    text = _ocr_image(file_bytes)
    if not text.strip():
        return []
    return [{"text": text.strip(), "page": 1, "order": 0, "type": "text"}]


def _ocr_image(image_bytes: bytes) -> str:
    try:
        from rapidocr_onnxruntime import RapidOCR
    except ImportError as exc:
        raise BusinessException("解析依赖未安装: rapidocr-onnxruntime") from exc
    engine = RapidOCR()
    result, _ = engine(image_bytes)
    if not result:
        return ""
    return "\n".join(item[1] for item in result)


# ---------------------------------------------------------------------------
# PDF 解析（fast / hi_res / ocr_only / auto）
# ---------------------------------------------------------------------------
def _parse_pdf(data: bytes, strategy: str) -> list[dict]:
    try:
        import pymupdf
    except ImportError as exc:
        raise BusinessException("解析依赖未安装: pymupdf") from exc

    doc = pymupdf.open(stream=data, filetype="pdf")
    try:
        if strategy == "hi_res":
            return _parse_pdf_hires(doc)
        if strategy == "ocr_only":
            return _parse_pdf_ocr(doc)
        blocks = _parse_pdf_fast(doc)
        if strategy == "auto" and not _blocks_has_text(blocks):
            blocks = _parse_pdf_ocr(doc)
        return blocks
    finally:
        doc.close()


def _parse_pdf_fast(doc) -> list[dict]:
    blocks = []
    order = 0
    for page_no in range(doc.page_count):
        text = doc[page_no].get_text().strip()
        if text:
            blocks.append({"text": text, "page": page_no + 1, "order": order, "type": "text"})
            order += 1
    return blocks


def _parse_pdf_ocr(doc) -> list[dict]:
    blocks = []
    order = 0
    for page_no in range(doc.page_count):
        pix = doc[page_no].get_pixmap(dpi=200)
        text = _ocr_image(pix.tobytes("png")).strip()
        if text:
            blocks.append({"text": text, "page": page_no + 1, "order": order, "type": "text"})
            order += 1
    return blocks


def _parse_pdf_hires(doc) -> list[dict]:
    """版面感知解析：文本/表格/公式/图片区域分类，每块保留 page/order 元数据。"""
    blocks = []
    order = 0
    for page_no in range(doc.page_count):
        page = doc[page_no]
        # 文本与公式区域（LaTeX 以 $$..$$ 原文保留在文本中）
        layout = page.get_text("dict")
        image_idx = 0
        for block in layout.get("blocks", []):
            # pymupdf block type：0=文本，1=图片
            block_type = block["type"]
            if block_type == 0:
                text = "".join(
                    span["text"] for line in block["lines"] for span in line["spans"]
                ).strip()
                if text:
                    blocks.append(
                        {"text": text, "page": page_no + 1, "order": order, "type": "text"}
                    )
                    order += 1
            elif block_type == 1:
                image_idx += 1
                blocks.append(
                    {
                        "text": f"[Image: 第{page_no + 1}页 图片{image_idx}]",
                        "page": page_no + 1,
                        "order": order,
                        "type": "image",
                    }
                )
                order += 1
        # 表格区域（Markdown 表格文本 + 行列元数据）
        try:
            tables = page.find_tables().tables
        except Exception:
            tables = []
        for table in tables:
            rows = table.extract() or []
            if not rows:
                continue
            md, rows_count, cols_count = _table_rows_to_markdown(rows)
            blocks.append(
                {
                    "text": md,
                    "page": page_no + 1,
                    "order": order,
                    "type": "table",
                    "table_rows": rows_count,
                    "table_cols": cols_count,
                }
            )
            order += 1
    return blocks


# ---------------------------------------------------------------------------
# Office / 文本解析（docx / xlsx / pptx / md / txt）
# ---------------------------------------------------------------------------
def _parse_office(data: bytes, ext: str) -> list[dict]:
    if ext == ".docx":
        return _parse_docx(data)
    if ext == ".xlsx":
        return _parse_xlsx(data)
    if ext == ".pptx":
        return _parse_pptx(data)
    return _parse_text(data)


def _parse_docx(data: bytes) -> list[dict]:
    try:
        import docx  # python-docx
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:
        raise BusinessException("解析依赖未安装: python-docx") from exc

    document = docx.Document(io.BytesIO(data))
    blocks = []
    order = 0
    # 按 body 子元素顺序遍历，保证段落与表格的原始阅读顺序
    for child in document.element.body.iterchildren():
        if child.tag == qn("w:p"):
            text = Paragraph(child, document).text.strip()
            if text:
                blocks.append({"text": text, "page": 1, "order": order, "type": "text"})
                order += 1
        elif child.tag == qn("w:tbl"):
            table = Table(child, document)
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            md, rows_count, cols_count = _table_rows_to_markdown(rows)
            blocks.append(
                {
                    "text": md,
                    "page": 1,
                    "order": order,
                    "type": "table",
                    "table_rows": rows_count,
                    "table_cols": cols_count,
                }
            )
            order += 1
    return blocks


def _parse_xlsx(data: bytes) -> list[dict]:
    try:
        import openpyxl
    except ImportError as exc:
        raise BusinessException("解析依赖未安装: openpyxl") from exc

    workbook = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    blocks = []
    order = 0
    for sheet in workbook.worksheets:
        # 合并单元格：首格值填充到合并区域其余格
        merged_map = {}
        for rng in sheet.merged_cells.ranges:
            top_left = sheet.cell(rng.min_row, rng.min_col).value
            for row in range(rng.min_row, rng.max_row + 1):
                for col in range(rng.min_col, rng.max_col + 1):
                    merged_map[(row, col)] = top_left
        rows_iter = sheet.iter_rows()
        header_cells = next(rows_iter, [])
        headers = [
            str(cell.value).strip() if cell.value is not None else f"列{i + 1}"
            for i, cell in enumerate(header_cells)
        ]
        for row in rows_iter:
            items = []
            for i, cell in enumerate(row):
                if i >= len(headers):
                    break
                value = cell.value
                if value is None:
                    value = merged_map.get((cell.row, cell.column))
                if value is None:
                    continue
                items.append(f"{headers[i]}: {str(value).strip()}")
            if items:
                blocks.append(
                    {
                        "text": ", ".join(items),
                        "page": 1,
                        "order": order,
                        "type": "table",
                        "table_rows": 1,  # 每行一个数据块
                        "table_cols": len(headers),
                        "sheet": sheet.title,
                    }
                )
                order += 1
    return blocks


def _parse_pptx(data: bytes) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise BusinessException("解析依赖未安装: python-pptx") from exc

    prs = Presentation(io.BytesIO(data))
    blocks = []
    order = 0
    for slide_no, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        texts.append(text)
            if getattr(shape, "has_table", False) and shape.has_table:
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                md, _, _ = _table_rows_to_markdown(rows)
                texts.append(md)
        combined = "\n".join(texts)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                combined = f"{combined}\n[备注] {notes}" if combined else f"[备注] {notes}"
        if combined.strip():
            blocks.append(
                {"text": combined.strip(), "page": slide_no, "order": order, "type": "text"}
            )
            order += 1
    return blocks


def _parse_text(data: bytes) -> list[dict]:
    text = _decode_text(data)
    if not text.strip():
        return []
    return [{"text": text.strip(), "page": 1, "order": 0, "type": "text"}]


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8", "gbk"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


# ---------------------------------------------------------------------------
# 表格工具
# ---------------------------------------------------------------------------
def _table_rows_to_markdown(rows: list[list[str]]) -> tuple[str, int, int]:
    """将二维行转 Markdown 表格，返回 (markdown, 数据行数, 列数)。"""
    if not rows:
        return "", 0, 0
    headers = rows[0]
    cols = len(headers)
    data_rows = len(rows) - 1
    lines = ["| " + " | ".join(_cell(h) for h in headers) + " |"]
    lines.append("|" + "|".join("---" for _ in headers) + "|")
    for row in rows[1:]:
        cells = [_cell(c) for c in row]
        if len(cells) < cols:
            cells += [""] * (cols - len(cells))
        lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines), data_rows, cols


def _cell(value: str) -> str:
    return (value or "").replace("\n", " ").replace("|", "\\|")
