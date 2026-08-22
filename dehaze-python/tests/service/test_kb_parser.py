import io

import pytest

from app.core.code import ResultCode
from app.core.exceptions import BusinessException
from app.service.kb.chunking_engine import chunk_text
from app.service.kb.document_parser import parse_document, parse_html


def _make_text(text: str) -> bytes:
    return text.encode("utf-8")


def _make_docx(paragraphs: list[str], table_rows: list[list[str]]) -> bytes:
    from docx import Document

    doc = Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    if table_rows:
        t = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        for r, row in enumerate(table_rows):
            for c, val in enumerate(row):
                t.cell(r, c).text = val
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx(headers: list, rows: list[list], merged: bool = False) -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.append(headers)
    for row in rows:
        ws.append(row)
    if merged and len(rows) >= 2:
        ws.merge_cells("A2:A3")
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(slides_data: list[dict]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for sd in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        texts = sd.get("texts", [])
        if texts:
            tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5)).text_frame
            tf.paragraphs[0].text = texts[0]
            for text in texts[1:]:
                tf.add_paragraph().text = text
        if sd.get("notes"):
            slide.notes_slide.notes_text_frame.text = sd["notes"]
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


_CJK_FONT = "/usr/share/fonts/google-noto-cjk/NotoSansCJKsc-Regular.otf"


def _make_pdf(text: str) -> bytes:
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), text, fontname="NotoSansCJKsc", fontfile=_CJK_FONT)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pdf_blocks(text_blocks: list[str]) -> bytes:
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    y = 50
    for i, text in enumerate(text_blocks):
        page.insert_text((72, y), f"{i}: {text}", fontname="NotoSansCJKsc", fontfile=_CJK_FONT)
        y += 30
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _corrupt_bytes(seed: int, size: int = 256) -> bytes:
    import random

    rng = random.Random(seed)
    return bytes(rng.randint(0, 255) for _ in range(size))


class TestTextFormats:
    def test_md_parsed_with_chinese(self):
        doc = parse_document(
            _make_text("# 系统上线说明\n\n本文档介绍去雾系统的部署步骤与回滚方案"), "readme.md"
        )
        assert "系统上线说明" in doc.content
        assert "部署步骤" in doc.content and "回滚方案" in doc.content

    def test_txt_parsed_with_chinese(self):
        doc = parse_document(_make_text("安装依赖：pip install -r requirements.txt"), "note.txt")
        assert doc.content == "安装依赖：pip install -r requirements.txt"

    def test_txt_gbk_fallback(self):
        data = "导出报表（GBK编码）".encode("gbk")
        doc = parse_document(data, "legacy.txt")
        assert "导出报表" in doc.content


class TestDirtyTextCorpus:
    def test_fullwidth_halfwidth_punctuation_mixed(self):
        text = "注意：config.yaml, 请核对（v2）。\n系统提示: 重试 3 次。\nQ：可用吗？A：可以。"
        doc = parse_document(_make_text(text), "mixed.txt")
        assert "config.yaml" in doc.content
        assert "重试 3 次" in doc.content
        assert "可用吗" in doc.content

    def test_crlf_lf_mixed(self):
        text = "行一\r\n行二\n行三\r\n行四"
        doc = parse_document(_make_text(text), "crlf.txt")
        for line in ("行一", "行二", "行三", "行四"):
            assert line in doc.content

    def test_utf8_bom_prefix(self):
        text = "\ufeff故障告警：CPU 使用率超过 90%"
        doc = parse_document(text.encode("utf-8"), "bom.txt")
        assert "故障告警" in doc.content

    def test_zero_width_chars(self):
        text = "\u200b系统启动\u200b完成\u200b"
        doc = parse_document(_make_text(text), "zw.txt")
        assert "系统启动" in doc.content and "完成" in doc.content

    def test_long_unbroken_line(self):
        line = "abcdefghABCDEFGH0123456789" * 400
        doc = parse_document(_make_text(line), "long.txt")
        assert len(doc.content) == len(line)
        assert line[:16] in doc.content

    def test_consecutive_blank_lines(self):
        text = "段落甲\n\n\n\n段落乙\n\n\n段落丁"
        doc = parse_document(_make_text(text), "blank.txt")
        assert "段落甲" in doc.content
        assert "段落丁" in doc.content


class TestDocx:
    def test_paragraphs_and_table(self):
        data = _make_docx(
            ["第一章 项目背景", "系统采用微服务架构部署"],
            [["模块", "负责人"], ["解析引擎", "张三"]],
        )
        doc = parse_document(data, "doc.docx")
        assert "第一章 项目背景" in doc.content
        assert "微服务架构" in doc.content
        assert "模块" in doc.content and "张三" in doc.content

    def test_table_block_metadata(self):
        data = _make_docx(["说明"], [["名称", "数量"], ["苹果", "5"]])
        doc = parse_document(data, "doc.docx")
        table_blocks = [b for b in doc.blocks if b["type"] == "table"]
        assert len(table_blocks) == 1
        assert table_blocks[0]["table_rows"] == 1
        assert table_blocks[0]["table_cols"] == 2


class TestTableVariants:
    def test_docx_header_only_table(self):
        data = _make_docx([], [["名称", "数量"]])
        doc = parse_document(data, "doc.docx")
        block = doc.blocks[0]
        assert block["type"] == "table"
        assert block["table_rows"] == 0
        assert block["table_cols"] == 2
        assert "| 名称 | 数量 |" in doc.content

    def test_docx_table_without_header(self):
        data = _make_docx([], [["苹果", "5"], ["香蕉", "3"]])
        doc = parse_document(data, "doc.docx")
        assert doc.blocks[0]["table_rows"] == 1
        assert doc.blocks[0]["table_cols"] == 2
        assert "苹果" in doc.content and "香蕉" in doc.content

    def test_docx_cell_pipe_escaped(self):
        data = _make_docx([], [["名称", "数量"], ["a|b", "5"]])
        doc = parse_document(data, "doc.docx")
        assert "a\\|b" in doc.content

    def test_xlsx_header_only_sheet_empty(self):
        data = _make_xlsx(["名称", "数量"], [])
        doc = parse_document(data, "sheet.xlsx")
        assert doc.content == ""

    def test_xlsx_cell_pipe_preserved(self):
        data = _make_xlsx(["名称", "数量"], [["a|b", "5"]])
        doc = parse_document(data, "sheet.xlsx")
        assert "名称: a|b" in doc.content


class TestXlsx:
    def test_structured_text_and_merge_fill(self):
        data = _make_xlsx(["姓名", "城市"], [["张三", "北京"], ["李四", "上海"]], merged=True)
        doc = parse_document(data, "sheet.xlsx")
        assert "姓名: 张三, 城市: 北京" in doc.content
        assert "姓名: 张三, 城市: 上海" in doc.content

    def test_table_block_metadata_unified_with_docx(self):
        data = _make_xlsx(["项目", "预算", "备注"], [["甲", "100", "x"], ["乙", "200", "y"]])
        doc = parse_document(data, "sheet.xlsx")
        assert all(b["type"] == "table" for b in doc.blocks)
        assert all(b["table_rows"] == 1 for b in doc.blocks)
        assert all(b["table_cols"] == 3 for b in doc.blocks)


class TestPptx:
    def test_slide_text_and_notes(self):
        data = _make_pptx([{"texts": ["季度总结", "营收环比 +8%"], "notes": "数据来自财报"}])
        doc = parse_document(data, "deck.pptx")
        assert "季度总结" in doc.content
        assert "营收环比 +8%" in doc.content
        assert "数据来自财报" in doc.content
        assert doc.blocks[0]["page"] == 1


class TestPdf:
    def test_fast_strategy_text_extracted(self):
        data = _make_pdf("本期财务摘要：营收同比增长12%")
        doc = parse_document(data, "file.pdf", strategy="fast")
        assert "营收" in doc.content and "12%" in doc.content
        assert doc.blocks[0]["type"] == "text"

    def test_hires_blocks_metadata(self):
        data = _make_pdf_blocks(["第一块文本", "第二块文本", "第三块文本"])
        doc = parse_document(data, "file.pdf", strategy="hi_res")
        assert len(doc.blocks) >= 1
        first = doc.blocks[0]
        assert first["type"] == "text"
        assert "page" in first and "order" in first
        assert first["order"] == 0

    def test_ocr_only_strategy_extracts_text(self):
        data = _make_pdf("扫描页")
        doc = parse_document(data, "file.pdf", strategy="ocr_only")
        assert "扫描页" in doc.content
        assert doc.blocks[0]["type"] == "text"
        assert doc.blocks[0]["page"] == 1


class TestUnsupportedFormat:
    def test_zip_raises_business_error(self):
        with pytest.raises(BusinessException) as exc:
            parse_document(b"PK\x03\x04", "archive.zip")
        assert exc.value.code == ResultCode.USER_UPLOAD_FILE_TYPE_NOT_MATCH


class TestParseHtml:
    def test_noise_tags_removed(self):
        html = """
        <html><body>
        <script>alert(1)</script>
        <nav>导航栏</nav>
        <style>.x{}</style>
        <p>正文内容</p>
        </body></html>
        """
        text = parse_html(html)
        assert "正文内容" in text
        assert "导航栏" not in text
        assert "alert" not in text

    def test_no_tag_residue(self):
        html = "<div><p>第一行</p><p>第二行</p></div>"
        text = parse_html(html)
        assert "<" not in text and ">" not in text
        assert "第一行" in text and "第二行" in text

    def test_script_content_with_chinese_and_table_not_in_body(self):
        html = """<html><body>
        <script>var x = "中文&<table>";</script>
        <p>真实正文</p>
        </body></html>"""
        text = parse_html(html)
        assert "真实正文" in text
        assert "var x" not in text

    def test_html_entities_decoded(self):
        html = "<p>Tom &amp; Jerry &lt;b&gt;bold&lt;/b&gt; &nbsp; end</p>"
        text = parse_html(html)
        assert "&" in text and "Tom" in text and "Jerry" in text
        assert "<b>" in text and "</b>" in text

    def test_body_less_malformed_html(self):
        html = "裸文本行一\n裸文本行二"
        assert parse_html(html).strip() == "裸文本行一\n裸文本行二".strip()

    def test_nested_tables(self):
        html = "<table><tr><td>外层<td>内层</td></tr></table>"
        text = parse_html(html)
        assert "外层" in text and "内层" in text


class TestEmptyFiles:
    @pytest.mark.parametrize("ext", [".pdf", ".docx", ".xlsx", ".pptx"])
    def test_zero_byte_binary_returns_empty(self, ext):
        doc = parse_document(b"", f"empty{ext}")
        assert doc.content == "" and doc.blocks == []

    @pytest.mark.parametrize("ext", [".txt", ".md"])
    def test_zero_byte_text_returns_empty(self, ext):
        doc = parse_document(b"", f"empty{ext}")
        assert doc.content == ""

    @pytest.mark.parametrize("ext", ["pdf", "docx", "xlsx", "pptx"])
    def test_corrupt_bytes_returns_business_error_or_empty(self, ext):
        data = _corrupt_bytes(12345)
        try:
            doc = parse_document(data, f"corrupt.{ext}")
            assert doc.content == ""
        except BusinessException as exc:
            assert exc.code == ResultCode.BUSINESS_ERROR

    def test_disguised_pdf_magic_as_txt(self):
        data = b"%PDF-1.4\n%\xe2\xe3\xcf\xd3"
        doc = parse_document(data, "fake.txt")
        assert "%PDF-1.4" in doc.content

    def test_disguised_text_as_pdf_returns_business_error(self):
        with pytest.raises(BusinessException):
            parse_document(b"This is plain text but disguised as PDF.", "fake.pdf", strategy="fast")


class TestXlsxAdversarial:
    def test_empty_sheet_returns_empty(self):
        doc = parse_document(_make_xlsx([], []), "empty.xlsx")
        assert doc.content == ""

    def test_sheet_with_only_blank_cells(self):
        doc = parse_document(_make_xlsx([None, None], [[None, None]]), "blank.xlsx")
        assert doc.content == ""

    def test_data_interspersed_with_many_blank_rows(self):
        blank_rows = [[None, None]] * 20
        doc = parse_document(_make_xlsx(["项目", "值"], blank_rows + [["甲", "100"]]), "sparse.xlsx")
        assert "项目: 甲, 值: 100" in doc.content


class TestDocxPptxAdversarial:
    def test_empty_docx_returns_empty(self):
        parsed = parse_document(_make_docx([], []), "empty.docx")
        assert parsed.content == ""

    def test_pptx_with_no_text_only_empty_slides(self):
        parsed = parse_document(_make_pptx([{}, {}]), "empty.pptx")
        assert parsed.content == ""


class TestChunkVariantRecognition:
    def test_header_only_table_rows_zero(self):
        table = "| 名称 | 数量 |\n| --- | --- |"
        t = next(c for c in chunk_text(table, "table", 800, 0) if c.metadata.get("type") == "table")
        assert t.metadata["rows"] == 0 and t.metadata["cols"] == 2

    def test_data_only_table_first_row_as_header(self):
        table = "| 苹果 | 5 |\n| 香蕉 | 3 |"
        t = next(c for c in chunk_text(table, "table", 800, 0) if c.metadata.get("type") == "table")
        assert t.metadata["rows"] == 1 and t.metadata["cols"] == 2

    def test_escaped_pipe_cell(self):
        table = "| 名称 | 数量 |\n| --- | --- |\n| a\\|b | 5 |"
        t = next(c for c in chunk_text(table, "table", 800, 0) if c.metadata.get("type") == "table")
        assert t.metadata["rows"] == 1 and t.metadata["cols"] == 2
        assert "a\\|b" in t.content

    def test_mixed_fullwidth_halfwidth_qa_colon(self):
        text = "Q：一\nA: 一答\nQ: 二\nA：二答"
        chunks = chunk_text(text, "qa", 800, 0)
        assert [c.metadata["type"] for c in chunks] == ["question", "answer", "question", "answer"]
        assert [c.content for c in chunks] == ["一", "一答", "二", "二答"]

    def test_tilde_fenced_pipe_not_table(self):
        md = "说明\n~~~\n| a | b |\n| --- | --- |\n~~~\n结尾"
        chunks = chunk_text(md, "table", 800, 0)
        assert all(c.metadata.get("type") != "table" for c in chunks)


class TestPerformanceSmoke:
    def test_large_chinese_fixed_chunking_under_5s(self):
        import time

        para = (
            "去雾系统 v2.3 于 2025-06-01 发布，新增批量解析与语义分块能力；"
            "支持 PDF/Word/Excel/PPT 与常见文本格式，兼容 GBK 与 UTF-8 编码，"
            "并内置 RapidOCR 图文识别，适用于客服工单与研发文档等场景。"
        )
        text = para * 1200
        start = time.perf_counter()
        chunks = chunk_text(text, "fixed", chunk_size=500, chunk_overlap=50)
        elapsed = time.perf_counter() - start
        assert chunks
        assert elapsed < 5.0, f"分块耗时 {elapsed:.2f}s 超过 5s 阈值"
